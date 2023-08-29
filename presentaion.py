import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a new PowerPoint presentation
presentation = Presentation()

# Title Slide
slide = presentation.slides.add_slide(presentation.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Salivary Gland Tumors"
subtitle.text = "Understanding Salivary Gland Tumors\nYour Name/Institution\nDate"

# Content Slides
slide_titles = [
    "Introduction",
    "Types of Salivary Glands",
    "What Are Salivary Gland Tumors?",
    "Classification of Salivary Gland Tumors",
    "Benign Salivary Gland Tumors",
    "Malignant Salivary Gland Tumors",
    "Causes and Risk Factors",
    "Symptoms and Diagnosis",
    "Treatment Options",
    "Prognosis",
    "Prevention and Awareness",
    "Case Studies",
    "Conclusion",
    "References"
]

for title_text in slide_titles:
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    title = slide.shapes.title
    title.text = title_text

# Save the presentation
presentation.save("Salivary_Gland_Tumors_Presentation.pptx")